#!/usr/bin/env python3
"""TRISK Desktop Installation Summary — 6-slide deck for bank IT/InfoSec."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys

# ── 1in1000 Brand Colors ──
CORAL = RGBColor(0xE8, 0x5D, 0x4E)
CORAL_DK = RGBColor(0xC9, 0x42, 0x33)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x33, 0x33, 0x33)
MID = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xAA, 0xAA, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GREEN_OK = RGBColor(0x2E, 0x7D, 0x32)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, font_size=14,
             bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12,
                  color=DARK, spacing=Pt(4), font_name="Arial"):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox


def add_kv_card(slide, left, top, width, height, label, value,
                label_color=MID, value_color=BLACK, bg_color=WHITE):
    """A small card with label on top, value below."""
    card = add_shape(slide, left, top, width, height, fill_color=bg_color,
                     line_color=RGBColor(0xE0, 0xE0, 0xE0))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = label_color
    p.font.name = "Arial"
    p.space_after = Pt(2)

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(11)
    p2.font.bold = False
    p2.font.color.rgb = value_color
    p2.font.name = "Arial"
    return card


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def build_deck(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, WHITE)

    # Coral accent bar at top
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()

    # Title
    add_text(slide1, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
             "TRISK Desktop", font_size=44, bold=True, color=BLACK)
    add_text(slide1, Inches(1.2), Inches(3.0), Inches(10), Inches(0.8),
             "Installation Requirements Summary", font_size=28, color=CORAL)

    # Separator line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1.2), Inches(3.9), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = CORAL
    line.line.fill.background()

    add_text(slide1, Inches(1.2), Inches(4.3), Inches(8), Inches(0.5),
             "Prepared by 1in1000 for [Bank Name]", font_size=16, color=MID)
    add_text(slide1, Inches(1.2), Inches(4.9), Inches(8), Inches(0.5),
             "Climate Transition Risk Stress Testing  |  ICAAP / ILAAP  |  Board Reporting",
             font_size=12, color=LIGHT)

    # Classification footer
    add_text(slide1, Inches(1.2), Inches(6.5), Inches(10), Inches(0.3),
             "INTERNAL  -  For IT / Information Security / Risk Technology Teams",
             font_size=9, color=LIGHT, alignment=PP_ALIGN.LEFT)

    add_notes(slide1, """TITLE SLIDE

Introduce the purpose: we're here to walk through the installation requirements
for TRISK Desktop before your IT team deploys it.

Key framing: this is a self-hosted, air-gapped tool. No cloud, no telemetry,
no data leaves your infrastructure. We designed it for exactly this kind of
institutional deployment.

The full 18-page Installation Guide has been provided separately. This deck
is the 6-slide summary for your review.""")

    # ════════════════════════════════════════════
    # SLIDE 2: What is TRISK Desktop?
    # ════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, WHITE)

    bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = CORAL
    bar2.line.fill.background()

    add_text(slide2, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "What is TRISK Desktop?", font_size=28, bold=True, color=BLACK)

    # One-liner
    add_text(slide2, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Self-hosted climate risk stress testing. Zero runtime network. 100% on-premise.",
             font_size=14, color=MID)

    # KV cards — 2 rows of 4
    cards = [
        ("RUNTIME NETWORK", "None - fully\nair-gapped"),
        ("DATA RESIDENCY", "100% on-premise\nNo cloud, no telemetry"),
        ("AUTHENTICATION", "HTTP Basic Auth\nor institutional SSO"),
        ("ENCRYPTION", "TLS 1.2+ via\nCaddy reverse proxy"),
        ("CONTAINER", "Non-root, all caps\ndropped, seccomp"),
        ("PII REQUIRED", "None - column\nallowlisting on upload"),
        ("RUNTIME DEPS", "Zero external\nall assets vendored"),
        ("USE CASE", "ICAAP, ILAAP\nBoard reporting"),
    ]

    x_start = Inches(0.8)
    y_row1 = Inches(1.9)
    y_row2 = Inches(4.0)
    card_w = Inches(2.8)
    card_h = Inches(1.6)
    gap = Inches(0.2)

    for i, (label, value) in enumerate(cards):
        row = 0 if i < 4 else 1
        col = i % 4
        x = x_start + col * (card_w + gap)
        y = y_row1 if row == 0 else y_row2
        add_kv_card(slide2, x, y, card_w, card_h, label, value)

    # Accent on first card (network = none)
    # We already created it, let's add a colored indicator on the slide
    dot = slide2.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x_start + Inches(2.35), y_row1 + Inches(0.15),
                                   Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN_OK
    dot.line.fill.background()

    add_notes(slide2, """WHAT IS TRISK DESKTOP

Walk through the 8 key properties that matter for IT/InfoSec review:

1. RUNTIME NETWORK: Zero outbound connections after installation. The Docker
   network is set to internal:true — the container literally cannot reach
   the internet. This is the single most important fact for your InfoSec team.

2. DATA RESIDENCY: All portfolio data stays on the analyst's machine or your
   internal server. No cloud endpoints, no analytics, no telemetry. The app
   doesn't even resolve DNS at runtime.

3. AUTHENTICATION: Ships with HTTP Basic Auth via the Caddy reverse proxy.
   For production, you can swap this for your institutional SSO (Keycloak,
   Azure AD, Okta) at the proxy layer — no application changes needed.

4. ENCRYPTION: TLS termination at Caddy. You provide your institution's
   certificates. HSTS enforced with 2-year max-age.

5. CONTAINER: Runs as non-root user 'trisk', all Linux capabilities dropped,
   read-only filesystem, custom seccomp profile blocking ptrace and key
   management syscalls.

6. PII: The application does not require any personally identifiable information.
   Column allowlisting on upload strips any fields not in the approved schema.

7. RUNTIME DEPENDENCIES: All fonts, JavaScript, CSS, and scenario data are
   vendored inside the container. No CDN calls, no external resources.

8. USE CASE: Credit risk teams use this for ICAAP climate stress test sections,
   ILAAP scenario analysis, and board-level reporting on transition risk.""")

    # ════════════════════════════════════════════
    # SLIDE 3: Prerequisites & Approvals
    # ════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, WHITE)

    bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = CORAL
    bar3.line.fill.background()

    add_text(slide3, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Prerequisites & Approvals", font_size=28, bold=True, color=BLACK)

    add_text(slide3, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "What your institution needs before we begin installation",
             font_size=14, color=MID)

    # Three columns
    col_w = Inches(3.6)
    col_h = Inches(4.5)
    col_gap = Inches(0.3)
    col_y = Inches(1.8)

    # Column 1: Governance
    c1 = add_shape(slide3, Inches(0.8), col_y, col_w, col_h,
                   fill_color=BG_LIGHT)
    add_text(slide3, Inches(1.1), col_y + Inches(0.2), Inches(3), Inches(0.4),
             "GOVERNANCE", font_size=11, bold=True, color=CORAL)
    checklist1 = [
        ("[ ]  CAB / Change Request approved", False),
        ("[ ]  CISO / InfoSec review complete", False),
        ("[ ]  Architecture review (if shared infra)", False),
        ("[ ]  License review (GPL-3.0 / LGPL-3.0)", False),
        ("[ ]  Vendor due diligence on 1in1000", False),
    ]
    add_multiline(slide3, Inches(1.1), col_y + Inches(0.7), Inches(3.2), Inches(3.5),
                  checklist1, font_size=11, color=DARK, spacing=Pt(6))

    # Column 2: Infrastructure
    c2 = add_shape(slide3, Inches(0.8) + col_w + col_gap, col_y, col_w, col_h,
                   fill_color=BG_LIGHT)
    add_text(slide3, Inches(1.1) + col_w + col_gap, col_y + Inches(0.2),
             Inches(3), Inches(0.4),
             "INFRASTRUCTURE", font_size=11, bold=True, color=CORAL)
    checklist2 = [
        ("[ ]  Docker Desktop in software catalog", False),
        ("[ ]  4 GB RAM allocated to Docker", False),
        ("[ ]  2 CPU cores minimum", False),
        ("[ ]  5 GB disk space free", False),
        ("[ ]  Port 443 (prod) or 3838 (dev)", False),
    ]
    add_multiline(slide3, Inches(1.1) + col_w + col_gap, col_y + Inches(0.7),
                  Inches(3.2), Inches(3.5),
                  checklist2, font_size=11, color=DARK, spacing=Pt(6))

    # Column 3: Security
    c3 = add_shape(slide3, Inches(0.8) + 2 * (col_w + col_gap), col_y, col_w, col_h,
                   fill_color=BG_LIGHT)
    add_text(slide3, Inches(1.1) + 2 * (col_w + col_gap), col_y + Inches(0.2),
             Inches(3), Inches(0.4),
             "SECURITY", font_size=11, bold=True, color=CORAL)
    checklist3 = [
        ("[ ]  TLS certificates from institution PKI", False),
        ("[ ]  Auth credentials configured", False),
        ("[ ]  SIEM log forwarding agreed", False),
        ("[ ]  Vulnerability scan plan agreed", False),
        ("[ ]  Scenario data file obtained", False),
    ]
    add_multiline(slide3, Inches(1.1) + 2 * (col_w + col_gap), col_y + Inches(0.7),
                  Inches(3.2), Inches(3.5),
                  checklist3, font_size=11, color=DARK, spacing=Pt(6))

    add_notes(slide3, """PREREQUISITES & APPROVALS

Walk through the three pillars. This maps to Section 2 and Section 6 of the
full Installation Guide.

GOVERNANCE:
- Standard change if Docker is already approved; normal change if first-time
- InfoSec review: the key document is the Security Architecture (next slide)
  plus the Dockerfile, seccomp profile, and docker-compose.yml
- License review: trisk.model is GPL-3.0, trisk.analysis is LGPL-3.0. Both
  are permissible for internal use without redistribution obligations since
  the bank is not distributing the software externally.

INFRASTRUCTURE:
- Docker Desktop 24.0+ required (or Docker Engine on Linux servers)
- 4 GB RAM is the container memory limit; recommend 16 GB total on machine
- Port 443 for production (TLS via Caddy), 3838 for development/UAT
- Windows 10 Pro/Enterprise, macOS 12+, Ubuntu 22.04+, RHEL 8+

SECURITY:
- TLS certs: request from your PKI team with CN=localhost and SAN including
  127.0.0.1. Self-signed available for testing.
- SIEM: Docker log driver can switch from json-file to syslog for forwarding
- Scenario data: ~200MB CSV, SHA256 verified at build. Can be transferred
  via approved media for air-gapped environments.""")

    # ════════════════════════════════════════════
    # SLIDE 4: Security Architecture
    # ════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, WHITE)

    bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar4.fill.solid()
    bar4.fill.fore_color.rgb = CORAL
    bar4.line.fill.background()

    add_text(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Security Architecture", font_size=28, bold=True, color=BLACK)

    # Network diagram — 3 boxes with arrows
    box_y = Inches(1.5)
    box_h = Inches(2.0)
    box_w = Inches(3.0)

    # Box 1: Browser
    b1 = add_shape(slide4, Inches(0.8), box_y, box_w, box_h,
                   fill_color=BG_LIGHT, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    tf1 = b1.text_frame
    tf1.margin_left = Pt(12)
    tf1.margin_top = Pt(12)
    tf1.vertical_anchor = MSO_ANCHOR.TOP
    p = tf1.paragraphs[0]
    p.text = "ANALYST BROWSER"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = MID
    p.font.name = "Arial"
    p2 = tf1.add_paragraph()
    p2.text = "localhost only\nNo LAN exposure"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK
    p2.font.name = "Arial"

    # Arrow 1
    add_text(slide4, Inches(3.9), box_y + Inches(0.6), Inches(1.2), Inches(0.8),
             "HTTPS\n(443)", font_size=11, bold=True, color=CORAL,
             alignment=PP_ALIGN.CENTER)
    arr1 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(3.95), box_y + Inches(1.3),
                                    Inches(1.0), Inches(0.3))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = CORAL
    arr1.line.fill.background()

    # Box 2: Caddy
    b2 = add_shape(slide4, Inches(5.1), box_y, box_w, box_h,
                   fill_color=WHITE, line_color=CORAL)
    tf2 = b2.text_frame
    tf2.margin_left = Pt(12)
    tf2.margin_top = Pt(12)
    tf2.vertical_anchor = MSO_ANCHOR.TOP
    p = tf2.paragraphs[0]
    p.text = "CADDY PROXY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CORAL
    p.font.name = "Arial"
    p2 = tf2.add_paragraph()
    p2.text = "TLS termination\nBasic Auth / SSO\nSecurity headers"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK
    p2.font.name = "Arial"

    # Arrow 2
    add_text(slide4, Inches(8.2), box_y + Inches(0.6), Inches(1.2), Inches(0.8),
             "HTTP\n(3838)", font_size=11, bold=True, color=MID,
             alignment=PP_ALIGN.CENTER)
    arr2 = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(8.25), box_y + Inches(1.3),
                                    Inches(1.0), Inches(0.3))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = MID
    arr2.line.fill.background()

    # Box 3: TRISK App
    b3 = add_shape(slide4, Inches(9.4), box_y, box_w, box_h,
                   fill_color=RGBColor(0xFA, 0xF0, 0xEE),
                   line_color=CORAL_DK)
    tf3 = b3.text_frame
    tf3.margin_left = Pt(12)
    tf3.margin_top = Pt(12)
    tf3.vertical_anchor = MSO_ANCHOR.TOP
    p = tf3.paragraphs[0]
    p.text = "TRISK APP"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CORAL_DK
    p.font.name = "Arial"
    p2 = tf3.add_paragraph()
    p2.text = "Internal network\nNo internet access\nRead-only filesystem"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK
    p2.font.name = "Arial"

    # Bottom section — 3 hardening pillars
    pill_y = Inches(4.2)
    pill_h = Inches(2.8)
    pill_w = Inches(3.6)

    pillars = [
        ("CONTAINER HARDENING", [
            "Non-root user (trisk)",
            "cap_drop: ALL",
            "no-new-privileges: true",
            "Read-only root filesystem",
            "Custom seccomp profile",
            "PID limit: 256",
        ]),
        ("SUPPLY CHAIN", [
            "Base image pinned by SHA256",
            "CRAN packages: date-locked snapshot",
            "GitHub packages: exact commit SHAs",
            "Scenario data: SHA256 verified",
            "All fonts vendored locally",
            "No CDN calls at runtime",
        ]),
        ("SECURITY HEADERS", [
            "X-Frame-Options: DENY",
            "X-Content-Type-Options: nosniff",
            "HSTS: 2-year max-age",
            "Referrer-Policy: no-referrer",
            "Camera/mic/geo/payment: blocked",
            "Server header: removed",
        ]),
    ]

    for i, (title, items) in enumerate(pillars):
        x = Inches(0.8) + i * (pill_w + Inches(0.3))
        card = add_shape(slide4, x, pill_y, pill_w, pill_h, fill_color=BG_LIGHT)

        add_text(slide4, x + Inches(0.2), pill_y + Inches(0.15),
                 Inches(3.2), Inches(0.3),
                 title, font_size=10, bold=True, color=CORAL)

        lines = [(item, False) for item in items]
        add_multiline(slide4, x + Inches(0.2), pill_y + Inches(0.55),
                      Inches(3.2), Inches(2.0),
                      lines, font_size=10, color=DARK, spacing=Pt(3))

    add_notes(slide4, """SECURITY ARCHITECTURE

This is the slide your InfoSec team will spend the most time on.

NETWORK DIAGRAM:
- The analyst's browser connects to Caddy on localhost:443 via HTTPS
- Caddy terminates TLS, handles authentication, applies security headers
- Caddy forwards to the TRISK app on an internal Docker network (port 3838)
- The internal Docker network has 'internal: true' — meaning the TRISK
  container cannot make ANY outbound connections. It cannot reach the
  internet, DNS, or any other network resource.

CONTAINER HARDENING:
- Non-root: runs as user 'trisk' with nologin shell
- cap_drop ALL: no Linux capabilities granted whatsoever
- Read-only filesystem: the container cannot write to its own image
- Writable only: /tmp (tmpfs, 2GB, noexec) and /data/output (noexec, nosuid)
- Custom seccomp profile: blocks ptrace, perf_event_open, keyctl, add_key,
  request_key, personality, userfaultfd, process_vm_readv/writev
- Memory: 4GB hard cap with 4GB swap cap to prevent OOM-swap spiral
- PIDs: limited to 256 to prevent fork bombs

SUPPLY CHAIN:
- Base image: rocker/r-ver:4.4.1 pinned by SHA256 digest, not a floating tag
- CRAN packages: locked to Posit PPM snapshot from 2025-02-01 — every build
  gets identical package versions
- GitHub packages: trisk.model and trisk.analysis pinned to exact Git commit
  SHAs — no surprise updates
- Scenario data: SHA256 checksum verified at build time, build fails on mismatch

The full security documentation is in SECURITY.md in the repository.
Your team can review the Dockerfile, docker-compose.yml, Caddyfile, and
seccomp profile directly.""")

    # ════════════════════════════════════════════
    # SLIDE 5: Installation — 3 Options
    # ════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, WHITE)

    bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar5.fill.solid()
    bar5.fill.fore_color.rgb = CORAL
    bar5.line.fill.background()

    add_text(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Installation  -  3 Options", font_size=28, bold=True, color=BLACK)

    add_text(slide5, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Choose based on your network policy. All produce identical results.",
             font_size=14, color=MID)

    # Three option cards
    opt_w = Inches(3.6)
    opt_h = Inches(4.5)
    opt_y = Inches(1.8)
    opt_gap = Inches(0.3)

    options = [
        ("OPTION A", "Standard Docker", CORAL,
         "Internet available at build time",
         [
             "1. git clone the repository",
             "2. Download scenario data",
             "3. Configure .env file",
             "4. Generate auth password",
             "5. Place TLS certificates",
             "6. docker compose up --build",
             "7. Verify at https://localhost",
         ]),
        ("OPTION B", "Bank Artifactory", RGBColor(0x4A, 0x7F, 0xA5),
         "Internal registry, pre-scanned image",
         [
             "Override BASE_IMAGE in .env",
             "Point to your artifactory",
             "Or pull pre-built image from",
             "  ghcr.io, scan, and re-host",
             "docker compose up -d",
             "",
             "No internet needed on target",
         ]),
        ("OPTION C", "Air-Gapped", RGBColor(0x5C, 0x6B, 0x73),
         "Zero network on target machine",
         [
             "Build on staging machine",
             "docker save | gzip images",
             "Transfer via approved media",
             "docker load on target",
             "Configure .env + TLS certs",
             "docker compose up -d",
             "",
         ]),
    ]

    for i, (badge, title, accent, subtitle, steps) in enumerate(options):
        x = Inches(0.8) + i * (opt_w + opt_gap)

        # Card background
        card = add_shape(slide5, x, opt_y, opt_w, opt_h,
                         fill_color=WHITE, line_color=RGBColor(0xE0, 0xE0, 0xE0))

        # Accent bar at top of card
        accent_bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              x, opt_y, opt_w, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # Badge
        add_text(slide5, x + Inches(0.2), opt_y + Inches(0.25),
                 Inches(1.5), Inches(0.3),
                 badge, font_size=9, bold=True, color=accent)

        # Title
        add_text(slide5, x + Inches(0.2), opt_y + Inches(0.55),
                 Inches(3.2), Inches(0.4),
                 title, font_size=16, bold=True, color=BLACK)

        # Subtitle
        add_text(slide5, x + Inches(0.2), opt_y + Inches(1.0),
                 Inches(3.2), Inches(0.3),
                 subtitle, font_size=10, color=MID)

        # Steps
        step_lines = [(s, False) for s in steps if s]
        add_multiline(slide5, x + Inches(0.2), opt_y + Inches(1.5),
                      Inches(3.2), Inches(2.8),
                      step_lines, font_size=10, color=DARK, spacing=Pt(4))

    # Fallback note
    add_text(slide5, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             "Fallback: Bare-metal R installation available if Docker is not approved. See INSTALL.md.",
             font_size=10, color=LIGHT)

    add_notes(slide5, """INSTALLATION - 3 OPTIONS

OPTION A — STANDARD DOCKER (most common):
7-step process that takes about 15 minutes. Requires outbound HTTPS to three
endpoints during build only: packagemanager.posit.co (CRAN packages),
github.com (TRISK model packages), storage.googleapis.com (scenario data).
After the build completes, the container never contacts the internet again.

OPTION B — BANK ARTIFACTORY:
For institutions using JFrog Artifactory, Sonatype Nexus, or similar internal
registries. Two sub-options: either override the base image at build time to
pull from your artifactory, or pull our pre-built image from ghcr.io, run your
CVE scan, host it internally, and have analysts pull from there.

OPTION C — AIR-GAPPED:
For environments with zero internet access on the target machine. Build on a
connected staging machine, export images with docker save, transfer via
approved media (encrypted USB, SFTP, etc.), and docker load on the target.
The full project directory also needs to be transferred for the configuration
files, Caddyfile, and TLS certificate setup.

BARE-METAL FALLBACK:
If Docker is not approved at all in your environment, TRISK can run directly
on R 4.4.1. This loses the container hardening (seccomp, capability drop,
read-only fs) but the application itself still runs on localhost only with
no outbound network. Additional endpoint security should be applied per
your institutional policy. Full instructions in INSTALL.md.""")

    # ════════════════════════════════════════════
    # SLIDE 6: Post-Installation & Operations
    # ════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, WHITE)

    bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar6.fill.solid()
    bar6.fill.fore_color.rgb = CORAL
    bar6.line.fill.background()

    add_text(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Post-Installation & Operations", font_size=28, bold=True, color=BLACK)

    # Four quadrant cards
    q_w = Inches(5.6)
    q_h = Inches(2.5)
    q_gap = Inches(0.3)

    quads = [
        (Inches(0.8), Inches(1.3), "VERIFICATION", CORAL, [
            ("Functional: Load demo data, run analysis, export results", False),
            ("Security: No outbound network, non-root, read-only fs, seccomp", False),
            ("Sign-off table provided in Installation Guide (Section 12)", False),
        ]),
        (Inches(0.8) + q_w + q_gap, Inches(1.3), "SIEM INTEGRATION", RGBColor(0x4A, 0x7F, 0xA5), [
            ("Switch Docker log driver from json-file to syslog", False),
            ("Forward to your SIEM at tcp://siem.bank.com:514", False),
            ("Audit events: uploads, analysis runs, exports, sessions", False),
        ]),
        (Inches(0.8), Inches(1.3) + q_h + q_gap, "VULNERABILITY SCANNING", RGBColor(0x5C, 0x6B, 0x73), [
            ("Pre-deployment: trivy image scan, remediate HIGH/CRITICAL", False),
            ("Ongoing: monthly rebuild + scan for new CVEs", False),
            ("SBOM generation: syft or trivy in SPDX-JSON format", False),
        ]),
        (Inches(0.8) + q_w + q_gap, Inches(1.3) + q_h + q_gap, "MAINTENANCE", GREEN_OK, [
            ("Backup: .env + tls/ + data/output/ only (stateless app)", False),
            ("Recovery: restore config, docker compose up --build", False),
            ("Decommission: docker compose down -v, remove directory", False),
        ]),
    ]

    for x, y, title, accent, items in quads:
        card = add_shape(slide6, x, y, q_w, q_h,
                         fill_color=WHITE, line_color=RGBColor(0xE0, 0xE0, 0xE0))
        # Accent bar
        ab = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, q_w, Inches(0.06))
        ab.fill.solid()
        ab.fill.fore_color.rgb = accent
        ab.line.fill.background()

        add_text(slide6, x + Inches(0.25), y + Inches(0.2),
                 Inches(5), Inches(0.3),
                 title, font_size=10, bold=True, color=accent)

        add_multiline(slide6, x + Inches(0.25), y + Inches(0.65),
                      Inches(5.1), Inches(1.7),
                      items, font_size=11, color=DARK, spacing=Pt(6))

    # Footer
    add_text(slide6, Inches(0.8), Inches(6.7), Inches(11), Inches(0.3),
             "Full Installation Guide: TRISK_Desktop_Installation_Guide_v1.pdf  |  For support: contact your 1in1000 engagement lead",
             font_size=9, color=LIGHT)

    add_notes(slide6, """POST-INSTALLATION & OPERATIONS

VERIFICATION (first day):
- Functional: Load the demo dataset, configure a scenario, run the analysis,
  verify charts render, test Excel/CSV export
- Security: Run the 4 docker exec commands from the guide to verify:
  1. wget fails (no outbound network)
  2. whoami returns 'trisk' (non-root)
  3. touch /test fails (read-only filesystem)
  4. Seccomp shows filter mode (2)
- Have the analyst and IT sign the verification table in Section 12

SIEM INTEGRATION:
- Replace the json-file logging driver in docker-compose.yml with syslog
- Point to your institutional SIEM endpoint
- Caddy access logs include authentication username and request details
- Application logs are structured JSON covering uploads, analysis runs,
  exports, and session lifecycle

VULNERABILITY SCANNING:
- Before deployment: run trivy or grype against the built image
- Policy: remediate all HIGH and CRITICAL before production
- Monthly: rebuild with --no-cache to pick up base image patches, then rescan
- Generate SBOM with syft for regulatory compliance documentation

MAINTENANCE:
- The application is stateless — no database, no session storage
- Only config files (.env, Caddyfile, TLS certs) and output data need backup
- Recovery is a single command: docker compose up --build
- Decommissioning is clean: docker compose down -v removes everything
- Archive data/output/ per your data retention policy before decommission

Next steps: share the full Installation Guide with your IT team, begin the
approval process (Section 2), and schedule the installation window.""")

    prs.save(output_path)
    print(f"Deck saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "TRISK_Installation_Summary_Deck_v1.pptx"
    build_deck(out)
